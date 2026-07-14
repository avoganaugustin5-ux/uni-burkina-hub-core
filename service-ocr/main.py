# ==============================================================================
# service-ocr — Port 8003
# UniBurkina Hub — AVOGAN Koudjo Augustin Sandaogo — INE N01213620231
# Version 3.0 — PaddleOCR uniquement (Gemini supprimé) :
#   - [REMOVED] Gemini Vision, corriger_texte_llm, _build_prompt_vision, _detect_type_document
#   - [KEEP] PaddleOCR + OpenCV layout + génération multi-format
#   - [KEEP] Stockage MinIO, workflow GED, JWT, SSE, CORS
#   - [KEEP] Tous les endpoints : retranscrire, retranscrire-stream, generer, soumettre-ged, valider
#   - [UPDATED] Score confiance adapté PaddleOCR uniquement
# ==============================================================================

import os, io, uuid, asyncio, httpx, logging, mimetypes
from pathlib import Path
from datetime import datetime
from typing import Optional, List

from fastapi import (
    FastAPI, File, UploadFile, Form, Header,
    HTTPException, BackgroundTasks, Depends
)
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from dotenv import load_dotenv


# ── OCR — PaddleOCR (moteur principal) ────────────────────────────────────────
from PIL import Image, ImageEnhance, ImageFilter
import numpy as np

try:
    from paddleocr import PaddleOCR
    _paddle = PaddleOCR(use_angle_cls=True, lang="fr", show_log=False)
    PADDLE_AVAILABLE = True
except Exception:
    PADDLE_AVAILABLE = False

# ── OpenCV pour détection et extraction de figures ────────────────────────────
try:
    import cv2
    CV2_AVAILABLE = True
except ImportError:
    CV2_AVAILABLE = False


# ── Génération PDF ─────────────────────────────────────────────────────────────
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib import colors
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer,
    Image as RLImage, HRFlowable
)
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY

# ── Génération DOCX ────────────────────────────────────────────────────────────
try:
    from docx import Document as DocxDocument
    from docx.shared import Pt, Cm as DocxCm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# ── Génération XLSX ────────────────────────────────────────────────────────────
try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

# ── Génération PPTX ────────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptPt
    from pptx.dml.color import RGBColor as PptRGB
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# ── Base de données ────────────────────────────────────────────────────────────
from sqlalchemy import (
    create_engine, Column, Integer, String, Text,
    DateTime, Float, ForeignKey, Boolean
)
from sqlalchemy.dialects.postgresql import UUID as PG_UUID
from sqlalchemy.orm import declarative_base, sessionmaker, Session
from sqlalchemy.sql import func

load_dotenv()
logging.basicConfig(level=logging.INFO)
log = logging.getLogger("service-ocr")

# Gemini supprimé en v3.0 — PaddleOCR uniquement
GEMINI_AVAILABLE = False

# ==============================================================================
# CONFIG
# ==============================================================================
DATABASE_URL = os.getenv(
    "DATABASE_URL",
    "postgresql://uniburkina_admin:UTS_Burkina2025!@localhost:5432/uniburkina_db"
)
API_AUTH     = os.getenv("API_AUTH",     "http://localhost:8001")
API_GED      = os.getenv("API_GED",      "http://localhost:8002")
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "")   # Clé Groq — optionnelle, active la correction Vision

# ── Stockage MinIO — module dédié (optionnel) ─────────────────────────────────
# Si ocr_storage_minio est absent ou si MinIO est arrêté, fallback stockage local.
_MINIO_AVAILABLE = False
try:
    from ocr_storage_minio import (
        init_ocr_storage,
        sauvegarder_image_minio,
        uploader_fichier_genere,
        telecharger_fichier_genere,
        fichier_genere_existe,
        get_tmp_output_dir,
        get_tmp_upload_dir,
    )
    _MINIO_AVAILABLE = True
except ImportError:
    logging.getLogger("service-ocr").warning(
        "[startup] ocr_storage_minio introuvable — stockage LOCAL activé (fallback)."
    )

_BASE_TMP     = Path(os.getenv("OCR_TMP_DIR", "tmp_ocr"))
_UPLOAD_LOCAL = _BASE_TMP / "uploads"
_OUTPUT_LOCAL = _BASE_TMP / "outputs"

if _MINIO_AVAILABLE:
    UPLOAD_DIR = get_tmp_upload_dir()
    OUTPUT_DIR = get_tmp_output_dir()
else:
    UPLOAD_DIR = _UPLOAD_LOCAL
    OUTPUT_DIR = _OUTPUT_LOCAL

    def init_ocr_storage():
        logging.getLogger("service-ocr").info("[fallback] stockage local.")

    def sauvegarder_image_minio(upload, idx: int, session_id: str) -> Path:
        import shutil
        ext  = Path(upload.filename or "image.jpg").suffix or ".jpg"
        dest = UPLOAD_DIR / f"{session_id}_{idx}{ext}"
        upload.file.seek(0)
        with open(dest, "wb") as fh:
            shutil.copyfileobj(upload.file, fh)
        return dest

    def uploader_fichier_genere(chemin: Path) -> str:
        import shutil
        dest = OUTPUT_DIR / chemin.name
        if chemin.resolve() != dest.resolve():
            shutil.copy2(str(chemin), str(dest))
        return str(dest)

    def telecharger_fichier_genere(cle: str):
        p = Path(cle)
        if not p.exists():
            raise FileNotFoundError(f"Fichier local introuvable : {cle}")
        return p.read_bytes(), p.name

    def fichier_genere_existe(cle: str) -> bool:
        return Path(cle).exists()

UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# Formats de sortie supportés
FORMATS_SUPPORTES = ["pdf", "docx", "xlsx", "pptx", "txt", "md"]

# Types MIME images acceptés
IMAGE_MIMES = {
    "image/jpeg", "image/png", "image/tiff",
    "image/bmp", "image/webp", "image/gif"
}

# Rôles autorisés à déposer sans validation
ROLES_DEPOT_DIRECT = {"ADMIN", "SOUS_ADMIN", "CHEF_DEPARTEMENT", "PRESIDENT",
                      "DIRECTEUR", "EMPLOYE", "CABINET", "SG", "VP_EIP", "VP_RCU"}
ROLES_SOUMISSION   = {"DELEGUE", "ETUDIANT", "ENSEIGNANT"} | ROLES_DEPOT_DIRECT

# ==============================================================================
# BASE DE DONNÉES
# ==============================================================================
engine       = create_engine(DATABASE_URL, pool_pre_ping=True, pool_size=10, max_overflow=20)
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base         = declarative_base()


class DocumentOCR(Base):
    __tablename__ = "documents_ocr"

    id_ocr             = Column(Integer, primary_key=True, index=True)
    nom_fichier_orig   = Column(String(500))
    chemin_source      = Column(String(1000))   # JSON list des chemins images sources
    chemin_pdf         = Column(String(1000))   # chemin du fichier généré final
    format_sortie      = Column(String(20), default="pdf")
    texte_extrait      = Column(Text)           # texte OCR brut (toutes images concat)
    texte_corrige      = Column(Text)           # texte après correction utilisateur
    nb_pages           = Column(Integer, default=1)
    nb_images          = Column(Integer, default=1)
    taille_fichier     = Column(Float)
    statut             = Column(String(50), default="EN_ATTENTE")
    # EN_ATTENTE | RETRANSCRIT | CORRIGE | GENERE | SOUMIS | VALIDE | REFUSE | ECHEC
    langue             = Column(String(20), default="fra")
    message_erreur     = Column(Text)
    date_soumission    = Column(DateTime, default=func.now())
    date_traitement    = Column(DateTime)
    # Liaison GED
    id_doc_ged         = Column(Integer, nullable=True)
    # Métadonnées du soumetteur
    id_utilisateur     = Column(PG_UUID(as_uuid=True), nullable=True)  # UUID (etait Integer par erreur)
    role_soumetteur    = Column(String(50), nullable=True)
    depot_direct       = Column(Boolean, default=False)  # True = pas besoin de validation


Base.metadata.create_all(bind=engine)


def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

# ==============================================================================
# APPLICATION
# ==============================================================================
app = FastAPI(
    title="UniBurkina Hub — Service OCR v2",
    version="2.5.0",
    description="OCR multi-images, multi-format, PaddleOCR, stockage MinIO — v3.0"
)

@app.on_event("startup")
async def startup_event():
    """Initialise les buckets MinIO au démarrage du service (non bloquant)."""
    try:
        init_ocr_storage()
        log.info("[startup] MinIO OCR initialisé avec succès.")
    except Exception as e:
        log.warning(
            f"[startup] MinIO non disponible au démarrage : {e} — "
            "le service démarre quand même. Relancez MinIO puis redémarrez service-ocr."
        )

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:8000", "http://localhost:8001",
        "http://localhost:8002", "http://localhost:8003",
        "http://localhost:8004", "http://localhost:8005",
        "http://127.0.0.1:8000", "http://127.0.0.1:8001",
        "http://127.0.0.1:8002", "http://127.0.0.1:8003",
        "http://127.0.0.1:8004", "http://127.0.0.1:8005",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ==============================================================================
# HELPERS
# ==============================================================================

async def verifier_token(
    authorization: Optional[str] = None,
    cookie: Optional[str] = Header(None)
) -> Optional[dict]:
    """Vérifie le JWT via service-auth — supporte header Bearer ET cookie HTTPOnly."""
    # Extraire le token depuis le header Authorization ou le cookie
    token = None
    if authorization and authorization.startswith("Bearer "):
        token = authorization.split(" ", 1)[1]
    elif cookie:
        # Cookie: access_token=xxx; autres=...
        for part in cookie.split(";"):
            part = part.strip()
            if part.startswith("access_token="):
                token = part.split("=", 1)[1].strip()
                break

    if not token:
        return None

    try:
        async with httpx.AsyncClient(timeout=5.0) as client:
            r = await client.get(
                f"{API_AUTH}/auth/me",
                headers={
                    "Authorization": f"Bearer {token}",
                    "Cookie": f"access_token={token}",
                }
            )
            if r.status_code == 200:
                return r.json()
    except Exception:
        pass
    return None


def _uuid_ou_none(valeur):
    """Convertit une chaine UUID en objet uuid.UUID ; retourne None si invalide/absente
    (id_utilisateur vient de /auth/me sous forme de chaine JSON, la colonne attend un UUID)."""
    if not valeur:
        return None
    try:
        return uuid.UUID(str(valeur))
    except (ValueError, AttributeError):
        return None


def ameliorer_image(img: Image.Image) -> Image.Image:
    """Pipeline d'amélioration image avant OCR pour maximiser la précision."""
    # Convertir en niveaux de gris
    img = img.convert("L")
    # Augmenter le contraste
    img = ImageEnhance.Contrast(img).enhance(2.0)
    # Netteté
    img = img.filter(ImageFilter.SHARPEN)
    # Binarisation adaptative simulée : seuil Otsu simplifié
    import numpy as np
    arr = np.array(img)
    threshold = arr.mean()
    arr = (arr > threshold).astype(np.uint8) * 255
    return Image.fromarray(arr)


def ocr_image(chemin: Path, langue: str = "fra") -> str:
    """Extrait le texte d'une image via PaddleOCR."""
    global _paddle, PADDLE_AVAILABLE
    chemin = Path(chemin)
    # ── Garde-fou : moteur non chargé ────────────────────────────────────────
    if not PADDLE_AVAILABLE:
        log.error(
            "PaddleOCR non disponible dans ce processus — "
            "pip install paddlepaddle==3.3.1 paddleocr==3.5.0 "
            "et relancer uvicorn SANS --reload"
        )
        raise HTTPException(
            status_code=503,
            detail=(
                "Moteur OCR non disponible (PADDLE_AVAILABLE=False). "
                "Vérifiez l'installation de PaddleOCR et relancez uvicorn sans --reload."
            )
        )
    try:
        result = _paddle.ocr(str(chemin), cls=True)
        if result and result[0]:
            lignes = [
                item[1][0]
                for bloc in result if bloc
                for item in bloc
                if item and len(item) >= 2
            ]
            return "\n".join(lignes).strip()
        # Résultat vide mais pas d'erreur : image illisible ou blanche
        log.warning(f"PaddleOCR n'a extrait aucun texte de {chemin.name} — image vide ou illisible ?")
        return ""
    except HTTPException:
        raise  # Laisser remonter les 503 déjà levés
    except Exception as e:
        log.error(f"OCR échoué sur {chemin}: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"Erreur OCR sur '{chemin.name}' : {e}"
        )





def _score_confiance(texte: str, moteur: str, has_figures: bool) -> int:
    """
    Calcule un score de confiance OCR de 0 à 100 basé sur :
      - le moteur utilisé (paddleocr+layout > paddleocr > echec)
      - la longueur du texte extrait
      - la présence de figures détectées
    """
    if moteur == "echec" or not texte:
        return 0
    # Base selon moteur
    base = 70 if moteur == "paddleocr+layout" else 55 if moteur == "paddleocr" else 30
    # Bonus longueur (max +20)
    longueur = len(texte.strip())
    bonus_longueur = min(20, longueur // 50)
    # Bonus figures détectées (+10)
    bonus_figures = 10 if has_figures else 0
    return min(100, base + bonus_longueur + bonus_figures)


def _nettoyer_texte_ocr(texte: str) -> str:
    """
    Nettoie le texte brut extrait par PaddleOCR :
      - Supprime les lignes vides multiples
      - Corrige la confusion visuelle I' → l' (police sans-serif)
      - Normalise les espaces et les tirets
    """
    if not texte:
        return ""
    import re
    # Confusion visuelle sans-serif : I' → l'
    texte = re.sub(r"\bI'", "l'", texte)
    texte = re.sub(r"\bI'", "l'", texte)   # apostrophe typographique
    # Espaces multiples → simple
    texte = re.sub(r" {2,}", " ", texte)
    # Plus de 2 sauts de ligne consécutifs → 2 max
    texte = re.sub(r"\n{3,}", "\n\n", texte)
    # Tirets longs → tiret simple
    texte = texte.replace("\u2014", "-").replace("\u2013", "-")
    return texte.strip()


def extraire_regions(chemin: Path) -> dict:
    """
    Détecte les zones texte et figure dans une image via OpenCV.
    Retourne {"blocs": [{"type":"texte"|"figure","y":int,"texte":str,"chemin_figure":Path|None}]}
    triés par position Y (ordre naturel de lecture).
    """
    if not CV2_AVAILABLE:
        log.warning("OpenCV non disponible — mode texte pur")
        try:
            texte_fallback = ocr_image(chemin)
        except HTTPException:
            raise
        except Exception as e:
            log.error(f"Fallback OCR échoué : {e}")
            texte_fallback = ""
        return {"blocs": [{"type": "texte", "y": 0,
                            "texte": texte_fallback, "chemin_figure": None}]}
    import cv2
    img_cv = cv2.imread(str(chemin))
    if img_cv is None:
        try:
            texte_fallback = ocr_image(chemin)
        except HTTPException:
            raise
        except Exception as e:
            log.error(f"Fallback OCR après imread None : {e}")
            texte_fallback = ""
        return {"blocs": [{"type": "texte", "y": 0,
                            "texte": texte_fallback, "chemin_figure": None}]}

    H, W = img_cv.shape[:2]
    gris = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
    _, bin_inv = cv2.threshold(gris, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)

    # Grands blocs compacts → figures/schémas
    k_fig = cv2.getStructuringElement(cv2.MORPH_RECT, (max(1, W // 4), max(1, H // 12)))
    dilate_fig = cv2.dilate(bin_inv, k_fig, iterations=2)
    contours_fig, _ = cv2.findContours(dilate_fig, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    MIN_AIRE = W * H * 0.03
    zones_fig = []
    for c in contours_fig:
        x, y, w, h = cv2.boundingRect(c)
        aire = w * h
        ratio = w / max(h, 1)
        if aire > MIN_AIRE and 0.3 < ratio < 7 and w < W * 0.97 and h < H * 0.97 and h > 35:
            zones_fig.append((x, y, w, h))

    # Fusion des boîtes qui se chevauchent
    zones_fig.sort(key=lambda z: z[1])
    merged = []
    for z in zones_fig:
        x, y, w, h = z
        fused = False
        for i, (mx, my, mw, mh) in enumerate(merged):
            if y < my+mh+25 and y+h > my-25 and x < mx+mw+25 and x+w > mx-25:
                merged[i] = (min(mx,x), min(my,y),
                              max(mx+mw,x+w)-min(mx,x), max(my+mh,y+h)-min(my,y))
                fused = True
                break
        if not fused:
            merged.append(z)

    # Pas de figure → texte pur
    if not merged:
        return {"blocs": [{"type": "texte", "y": 0,
                            "texte": ocr_image(chemin), "chemin_figure": None}]}

    figures_dir = UPLOAD_DIR / "figures"
    figures_dir.mkdir(parents=True, exist_ok=True)
    blocs = []

    # Découper et sauvegarder chaque figure
    for i, (x, y, w, h) in enumerate(merged):
        mg = 10
        crop = img_cv[max(0,y-mg):min(H,y+h+mg), max(0,x-mg):min(W,x+w+mg)]
        dest_fig = figures_dir / f"{chemin.stem}_fig{i:02d}.png"
        cv2.imwrite(str(dest_fig), crop)
        blocs.append({"type": "figure", "y": y, "texte": "", "chemin_figure": dest_fig})

    # Blocs texte = contours hors figures
    k_txt = cv2.getStructuringElement(cv2.MORPH_RECT, (max(1, W // 8), 5))
    dilate_txt = cv2.dilate(bin_inv, k_txt, iterations=2)
    contours_txt, _ = cv2.findContours(dilate_txt, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    def dans_figure(x, y, w, h):
        for (fx, fy, fw, fh) in merged:
            if x >= fx-12 and y >= fy-12 and x+w <= fx+fw+12 and y+h <= fy+fh+12:
                return True
        return False

    bandes = []
    for c in contours_txt:
        x, y, w, h = cv2.boundingRect(c)
        if w < 20 or h < 5:
            continue
        if not dans_figure(x, y, w, h):
            bandes.append((x, y, w, h))

    bandes.sort(key=lambda b: b[1])
    if bandes:
        groupes = [[bandes[0]]]
        for bt in bandes[1:]:
            dernier = groupes[-1][-1]
            if bt[1] - (dernier[1] + dernier[3]) > 40:
                groupes.append([bt])
            else:
                groupes[-1].append(bt)
        for groupe in groupes:
            gx  = min(b[0] for b in groupe)
            gy  = min(b[1] for b in groupe)
            gx2 = max(b[0]+b[2] for b in groupe)
            gy2 = max(b[1]+b[3] for b in groupe)
            mg = 4
            crop_txt = img_cv[max(0,gy-mg):min(H,gy2+mg), max(0,gx-mg):min(W,gx2+mg)]
            tmp = figures_dir / f"{chemin.stem}_txt_{gy:04d}.png"
            cv2.imwrite(str(tmp), crop_txt)
            texte_bande = ocr_image(tmp)
            if texte_bande.strip():
                blocs.append({"type": "texte", "y": gy,
                               "texte": texte_bande, "chemin_figure": None})

    blocs.sort(key=lambda b: b["y"])
    nb_fig = sum(1 for b in blocs if b["type"] == "figure")
    nb_txt = sum(1 for b in blocs if b["type"] == "texte")
    log.info(f"Layout {chemin.name} → {nb_fig} figure(s), {nb_txt} bloc(s) texte")
    return {"blocs": blocs}


def sauvegarder_image(upload: UploadFile, idx: int, session_id: str) -> Path:
    """
    Sauvegarde une image uploadée.
    Délègue à ocr_storage_minio : écrit localement (pour Pillow/cv2)
    ET upload vers MinIO bucket uniburkina-ocr-uploads.
    Retourne le chemin local temporaire (compatible avec tout le code OCR existant).
    """
    return sauvegarder_image_minio(upload, idx, session_id)


# ==============================================================================
# GÉNÉRATEURS DE FORMAT
# ==============================================================================

def generer_pdf(texte: str, titre: str, images: List[Path] = None,
                mode_images: bool = False, blocs_layout: list = None) -> Path:
    """Génère un PDF depuis texte OCR ou directement depuis images."""
    nom = f"doc_{uuid.uuid4().hex[:10]}.pdf"
    dest = OUTPUT_DIR / nom
    doc  = SimpleDocTemplate(
        str(dest), pagesize=A4,
        leftMargin=2*cm, rightMargin=2*cm,
        topMargin=2.5*cm, bottomMargin=2*cm
    )
    styles = getSampleStyleSheet()
    titre_style = ParagraphStyle(
        "TitreUBH", parent=styles["Heading1"],
        fontSize=16, spaceAfter=12,
        textColor=colors.HexColor("#042C53"),
        alignment=TA_CENTER
    )
    corps_style = ParagraphStyle(
        "CorpsUBH", parent=styles["Normal"],
        fontSize=11, leading=16, spaceAfter=6,
        alignment=TA_JUSTIFY
    )
    contenu = []
    # En-tête
    contenu.append(Paragraph("UniBurkina Hub — Document généré par OCR", titre_style))
    contenu.append(HRFlowable(width="100%", thickness=1.5,
                               color=colors.HexColor("#C89A2E"), spaceAfter=12))
    if titre:
        contenu.append(Paragraph(titre, titre_style))
        contenu.append(Spacer(1, 0.5*cm))

    if mode_images and images:
        # Mode direct : insérer les images dans l'ordre
        for i, img_path in enumerate(images):
            try:
                with Image.open(img_path) as img:
                    w, h = img.size
                    max_w = 16 * cm
                    ratio = min(max_w / w, 22*cm / h)
                    rl_img = RLImage(str(img_path), width=w*ratio, height=h*ratio)
                    contenu.append(rl_img)
                    contenu.append(Spacer(1, 0.3*cm))
                    contenu.append(Paragraph(f"<i>Image {i+1}/{len(images)}</i>",
                                             styles["Italic"]))
                    contenu.append(Spacer(1, 0.5*cm))
            except Exception as e:
                log.warning(f"Image {img_path} ignorée : {e}")
    elif blocs_layout:
        # Mode layout : intercaler texte et figures dans l'ordre naturel de lecture
        for bloc in blocs_layout:
            if bloc["type"] == "figure" and bloc.get("chemin_figure"):
                fig_path = Path(bloc["chemin_figure"])
                if fig_path.exists():
                    try:
                        with Image.open(fig_path) as img:
                            w, h = img.size
                            max_w = 14 * cm
                            ratio = min(max_w / w, 10*cm / h, 1.0)
                            rl_img = RLImage(str(fig_path), width=w*ratio, height=h*ratio)
                            contenu.append(rl_img)
                            contenu.append(Spacer(1, 0.3*cm))
                    except Exception as e:
                        log.warning(f"Figure {fig_path} ignorée : {e}")
            elif bloc["type"] == "texte" and bloc.get("texte", "").strip():
                for ligne in bloc["texte"].split("\n"):
                    ligne = ligne.strip()
                    if not ligne:
                        contenu.append(Spacer(1, 0.2*cm))
                        continue
                    try:
                        contenu.append(Paragraph(ligne, corps_style))
                    except Exception:
                        contenu.append(Paragraph(
                            ligne.replace("<", "&lt;").replace(">", "&gt;"), corps_style))
    else:
        # Mode texte pur : paragraphes depuis le texte OCR (corrigé ou brut)
        for ligne in (texte or "").split("\n"):
            ligne = ligne.strip()
            if not ligne:
                contenu.append(Spacer(1, 0.2*cm))
                continue
            try:
                contenu.append(Paragraph(ligne, corps_style))
            except Exception:
                contenu.append(Paragraph(ligne.replace("<", "&lt;").replace(">", "&gt;"), corps_style))

    # Pied de page
    contenu.append(Spacer(1, 1*cm))
    contenu.append(HRFlowable(width="100%", thickness=0.5,
                               color=colors.grey, spaceAfter=6))
    contenu.append(Paragraph(
        f"Généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')} "
        f"— Université Thomas SANKARA — UFR Sciences et Technologies",
        ParagraphStyle("footer", parent=styles["Italic"], fontSize=8,
                       textColor=colors.grey, alignment=TA_CENTER)
    ))
    doc.build(contenu)
    return dest


def generer_docx(texte: str, titre: str) -> Path:
    """Génère un fichier DOCX depuis le texte OCR."""
    if not DOCX_AVAILABLE:
        raise HTTPException(503, "python-docx non installé sur ce serveur")
    nom  = f"doc_{uuid.uuid4().hex[:10]}.docx"
    dest = OUTPUT_DIR / nom
    document = DocxDocument()
    # Style titre
    h = document.add_heading(titre or "Document OCR — UniBurkina Hub", 0)
    h.runs[0].font.color.rgb = RGBColor(0x04, 0x2C, 0x53)
    document.add_paragraph()
    # Contenu
    for ligne in (texte or "").split("\n"):
        ligne = ligne.strip()
        if not ligne:
            document.add_paragraph()
            continue
        p = document.add_paragraph(ligne)
        p.paragraph_format.space_after = Pt(4)
    # Pied de page
    section = document.sections[0]
    footer  = section.footer
    footer.paragraphs[0].text = (
        f"UniBurkina Hub — Généré le {datetime.now().strftime('%d/%m/%Y')} "
        f"— Université Thomas SANKARA"
    )
    footer.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
    document.save(str(dest))
    return dest


def generer_xlsx(texte: str, titre: str) -> Path:
    """Génère un fichier Excel depuis le texte OCR (une ligne = une ligne)."""
    if not XLSX_AVAILABLE:
        raise HTTPException(503, "openpyxl non installé sur ce serveur")
    nom  = f"doc_{uuid.uuid4().hex[:10]}.xlsx"
    dest = OUTPUT_DIR / nom
    wb   = openpyxl.Workbook()
    ws   = wb.active
    ws.title = titre[:31] if titre else "OCR Document"
    # En-tête
    ws["A1"] = titre or "Document OCR — UniBurkina Hub"
    ws["A1"].font = Font(bold=True, size=14, color="042C53")
    ws["A1"].fill = PatternFill("solid", fgColor="E6F1FB")
    ws.merge_cells("A1:D1")
    ws["A1"].alignment = Alignment(horizontal="center")
    row = 3
    for ligne in (texte or "").split("\n"):
        ligne = ligne.strip()
        if ligne:
            ws.cell(row=row, column=1, value=ligne)
            row += 1
    ws.column_dimensions["A"].width = 100
    wb.save(str(dest))
    return dest


def generer_pptx(texte: str, titre: str) -> Path:
    """Génère une présentation PPTX depuis le texte OCR."""
    if not PPTX_AVAILABLE:
        raise HTTPException(503, "python-pptx non installé sur ce serveur")
    nom  = f"doc_{uuid.uuid4().hex[:10]}.pptx"
    dest = OUTPUT_DIR / nom
    prs  = Presentation()
    # Slide titre
    layout = prs.slide_layouts[0]
    slide  = prs.slides.add_slide(layout)
    slide.shapes.title.text = titre or "Document OCR"
    slide.placeholders[1].text = (
        f"UniBurkina Hub — Université Thomas SANKARA\n"
        f"Généré le {datetime.now().strftime('%d/%m/%Y')}"
    )
    # Découper le texte en slides (toutes les 15 lignes)
    lignes = [l.strip() for l in (texte or "").split("\n") if l.strip()]
    chunks = [lignes[i:i+15] for i in range(0, len(lignes), 15)]
    for i, chunk in enumerate(chunks):
        layout_c = prs.slide_layouts[1]
        slide_c  = prs.slides.add_slide(layout_c)
        slide_c.shapes.title.text = f"Page {i+1}"
        tf = slide_c.placeholders[1]
        tf.text = "\n".join(chunk)
    prs.save(str(dest))
    return dest


def generer_txt(texte: str, titre: str) -> Path:
    """Génère un fichier TXT depuis le texte OCR."""
    nom  = f"doc_{uuid.uuid4().hex[:10]}.txt"
    dest = OUTPUT_DIR / nom
    entete = (
        f"UniBurkina Hub — Document généré par OCR\n"
        f"Université Thomas SANKARA — UFR Sciences et Technologies\n"
        f"Titre : {titre or 'Sans titre'}\n"
        f"Date  : {datetime.now().strftime('%d/%m/%Y à %H:%M')}\n"
        f"{'=' * 60}\n\n"
    )
    dest.write_text(entete + (texte or ""), encoding="utf-8")
    return dest


def generer_markdown(texte: str, titre: str) -> Path:
    """Génère un fichier Markdown depuis le texte OCR."""
    nom  = f"doc_{uuid.uuid4().hex[:10]}.md"
    dest = OUTPUT_DIR / nom
    contenu = (
        f"# {titre or 'Document OCR — UniBurkina Hub'}\n\n"
        f"> Généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}  \n"
        f"> Université Thomas SANKARA — UFR Sciences et Technologies\n\n"
        f"---\n\n"
    )
    for ligne in (texte or "").split("\n"):
        ligne = ligne.strip()
        contenu += (ligne + "\n") if ligne else "\n"
    dest.write_text(contenu, encoding="utf-8")
    return dest


def dispatcher_generation(
    format_sortie: str, texte: str, titre: str,
    images: List[Path] = None, mode_images: bool = False,
    blocs_layout: list = None
) -> Path:
    """Dispatche vers le bon générateur selon le format demandé."""
    f = format_sortie.lower().strip()
    if f == "pdf":
        return generer_pdf(texte, titre, images, mode_images, blocs_layout=blocs_layout)
    elif f == "docx":
        return generer_docx(texte, titre)
    elif f in ("xlsx", "xls"):
        return generer_xlsx(texte, titre)
    elif f in ("pptx", "ppt"):
        return generer_pptx(texte, titre)
    elif f == "txt":
        return generer_txt(texte, titre)
    elif f in ("md", "markdown"):
        return generer_markdown(texte, titre)
    else:
        raise HTTPException(400, f"Format non supporté : {format_sortie}. "
                                 f"Formats disponibles : {', '.join(FORMATS_SUPPORTES)}")


# ==============================================================================
# NOTIFICATION (background task)
# ==============================================================================

async def notifier_utilisateur(
    id_utilisateur: int, titre_notif: str, message: str, token: str
):
    """Envoie une notification via service-auth."""
    try:
        async with httpx.AsyncClient(timeout=5.0) as client:
            await client.post(
                f"{API_AUTH}/auth/notifications",
                json={
                    "titre":       titre_notif,
                    "message":     message,
                    "user_id":     id_utilisateur,
                },
                headers={"Authorization": f"Bearer {token}"}
            )
    except Exception as e:
        log.warning(f"Notification échouée pour user {id_utilisateur} : {e}")


# ==============================================================================
# ENDPOINTS
# ==============================================================================

# ── 1. SCAN MULTI-IMAGES → Retranscription seule ──────────────────────────────
@app.post("/ocr/scan-multi", summary="Scanner plusieurs images et retourner le texte extrait")
async def scan_multi(
    files: List[UploadFile] = File(...),
    langue: str = Form("fra"),
    ordre: str  = Form(""),   # ex: "0,2,1" pour réordonner les images
    authorization: Optional[str] = Header(None),
    db: Session = Depends(get_db)
):
    """
    Scanne plusieurs images dans l'ordre spécifié.
    Retourne le texte extrait de chaque image + le texte consolidé.
    L'utilisateur peut ensuite corriger ce texte avant de générer le document.
    """
    user = await verifier_token(authorization)

    # Validation des fichiers
    for f in files:
        if f.content_type not in IMAGE_MIMES:
            raise HTTPException(400, f"{f.filename} : format non supporté. "
                                     f"Formats acceptés : JPEG, PNG, TIFF, BMP, WebP")
        if f.size and f.size > 20 * 1024 * 1024:  # 20 Mo max
            raise HTTPException(413, f"{f.filename} : fichier trop lourd (max 20 Mo)")

    session_id  = uuid.uuid4().hex[:12]
    chemins     = []

    # Sauvegarder dans l'ordre d'upload
    for i, f in enumerate(files):
        chemin = sauvegarder_image(f, i, session_id)
        chemins.append(chemin)

    # Réordonner si demandé
    if ordre.strip():
        try:
            indices = [int(x) for x in ordre.split(",")]
            assert len(indices) == len(chemins)
            chemins = [chemins[i] for i in indices]
        except Exception:
            pass  # Ignorer si l'ordre est invalide, garder l'ordre d'upload

    # ── OCR sur chaque image : Gemini Vision en priorité, PaddleOCR en fallback ─
    resultats    = []
    texte_global = ""
    moteurs_utilises = []
    has_figures_global = False

    for i, chemin in enumerate(chemins):
        texte_img   = ""
        has_figures = False
        moteur_img  = "paddleocr"

        # ── PaddleOCR + détection layout OpenCV ──────────────────────────────
        try:
            layout_result = extraire_regions(chemin)
            blocs_img     = layout_result["blocs"]
            has_figures   = any(b["type"] == "figure" for b in blocs_img)
            texte_paddle  = "\n".join(
                b["texte"] for b in blocs_img if b["type"] == "texte"
            )
            texte_img  = _nettoyer_texte_ocr(texte_paddle)
            moteur_img = "paddleocr+layout" if has_figures else "paddleocr"
            log.info(f"Image {i+1} — PaddleOCR OK ({len(texte_img)} chars)")
        except HTTPException as e:
            log.error(f"Image {i+1} — PaddleOCR erreur {e.status_code}: {e.detail}")
            texte_img  = ""
            moteur_img = "echec"

        moteurs_utilises.append(moteur_img)
        has_figures_global = has_figures_global or has_figures
        resultats.append({
            "index":         i,
            "nom_fichier":   chemin.name,
            "texte_extrait": texte_img,
            "nb_chars":      len(texte_img),
            "has_figures":   has_figures,
            "moteur":        moteur_img,
        })
        texte_global += texte_img + "\n\n"

    # ── Résumé moteur utilisé ─────────────────────────────────────────────────
    moteur_principal = (
        "paddleocr+layout" if all(m == "paddleocr+layout" for m in moteurs_utilises)
        else (moteurs_utilises[0] if moteurs_utilises else "paddleocr")
    )

    texte_final = texte_global.strip()

    # Persister en BDD — texte_extrait = brut, texte_corrige = résultat final
    doc = DocumentOCR(
        nom_fichier_orig  = files[0].filename if files else "scan",
        chemin_source     = str([str(c) for c in chemins]),
        texte_extrait     = texte_final,
        texte_corrige     = texte_final,
        nb_pages          = len(files),
        nb_images         = len(files),
        statut            = "RETRANSCRIT",
        langue            = langue,
        date_traitement   = datetime.now(),
        id_utilisateur    = _uuid_ou_none(user.get("id")) if user else None,
        role_soumetteur   = user.get("role") if user else None,
    )
    db.add(doc); db.commit(); db.refresh(doc)

    return {
        "id_ocr":              doc.id_ocr,
        "session_id":          session_id,
        "nb_images":           len(files),
        "langue":              langue,
        "statut":              "CORRIGE",
        "moteur_ocr":          moteur_principal,
        "has_figures":         has_figures_global,
        "resultats_par_image": resultats,
        "texte_consolide":     texte_final,
        "texte_corrige":       texte_final,
        "llm_correction":      True,
        "message": (
            f"{len(files)} image(s) traitée(s) par PaddleOCR"
            + (" (layout/figures détectés)" if has_figures_global else "")
            + "."
        )
    }


# ── 2. CORRIGER LE TEXTE OCR ──────────────────────────────────────────────────
@app.patch("/ocr/documents/{id_ocr}/corriger",
           summary="Enregistrer le texte corrigé par l'utilisateur")
async def corriger_texte(
    id_ocr: int,
    texte_corrige: str = Form(...),
    authorization: Optional[str] = Header(None),
    db: Session = Depends(get_db)
):
    """
    Enregistre le texte corrigé par l'utilisateur après retranscription.
    Passe le statut de RETRANSCRIT → CORRIGE.
    """
    doc = db.query(DocumentOCR).filter(DocumentOCR.id_ocr == id_ocr).first()
    if not doc:
        raise HTTPException(404, "Document OCR introuvable")

    doc.texte_corrige = texte_corrige
    doc.statut        = "CORRIGE"
    db.commit()

    return {"id_ocr": id_ocr, "statut": "CORRIGE", "nb_chars": len(texte_corrige)}


# ── 2b. CORRECTION GROQ VISION (bouton "Corriger avec IA" du frontend) ────────
@app.post("/ocr/documents/{id_ocr}/corriger-llm",
          summary="Correction IA via Groq Vision (llama-4-scout) — remplace Gemini")
async def corriger_llm_manuel(
    id_ocr: int,
    authorization: Optional[str] = Header(None),
    db: Session = Depends(get_db)
):
    """
    Réactivé en v3.1 : envoie l'image source + texte PaddleOCR à Groq Vision
    pour correction orthographique, formules mathématiques et amélioration structure.

    Requiert GROQ_API_KEY dans .env — gratuit sur groq.com.
    Si la clé est absente ou si Groq échoue, retourne le texte PaddleOCR inchangé.
    """
    from ocr_engine import groq_vision_corriger, GROQ_AVAILABLE

    doc = db.query(DocumentOCR).filter(DocumentOCR.id_ocr == id_ocr).first()
    if not doc:
        raise HTTPException(404, "Document OCR introuvable")

    if not GROQ_AVAILABLE:
        raise HTTPException(
            503,
            "GROQ_API_KEY non configurée dans service-ocr/.env. "
            "Ajoutez GROQ_API_KEY=gsk_... et redémarrez le service."
        )

    # Récupérer la première image source depuis la BDD
    import json as _json
    try:
        chemins_src = _json.loads(doc.chemin_source.replace("'", '"'))
        chemins_valides = [Path(c) for c in chemins_src if Path(c).exists()]
    except Exception:
        chemins_valides = []

    if not chemins_valides:
        raise HTTPException(
            404,
            "Images sources introuvables (fichiers temporaires supprimés ?). "
            "Relancez une retranscription et corrigez immédiatement après."
        )

    # Texte de base : corrigé manuellement en priorité, sinon brut PaddleOCR
    texte_base = doc.texte_corrige or doc.texte_extrait or ""

    # Appeler Groq Vision sur chaque image et consolider
    textes_corriges = []
    moteurs         = []
    score_bonus_tot = 0

    for chemin in chemins_valides:
        resultat = await groq_vision_corriger(chemin, texte_base)
        textes_corriges.append(resultat["texte_corrige"])
        moteurs.append(resultat["moteur"])
        score_bonus_tot += resultat["score_bonus"]

    texte_final = "\n\n".join(textes_corriges).strip()
    moteur_final = "paddleocr+groq" if any("groq" in m for m in moteurs) else "paddleocr"

    # Sauvegarder en BDD
    doc.texte_corrige = texte_final
    doc.statut        = "CORRIGE"
    db.commit()

    log.info(f"Groq correction id_ocr={id_ocr} — {len(texte_final)} chars, moteur={moteur_final}")

    return {
        "id_ocr":         id_ocr,
        "statut":         "CORRIGE",
        "moteur":         moteur_final,
        "nb_chars":       len(texte_final),
        "texte_corrige":  texte_final,
        "message":        "Texte corrigé par Groq Vision (llama-4-scout)." if moteur_final == "paddleocr+groq"
                          else "Groq n'a pas pu corriger — texte PaddleOCR conservé.",
    }


# ── 3. GÉNÉRER LE DOCUMENT AU FORMAT VOULU ────────────────────────────────────
@app.post("/ocr/generer",
          summary="Générer un document au format choisi depuis texte ou images")
async def generer_document(
    id_ocr:        int            = Form(...),
    format_sortie: str            = Form("pdf"),
    titre:         str            = Form(""),
    mode:          str            = Form("texte"),  # "texte" | "images"
    authorization: Optional[str]  = Header(None),
    db:            Session        = Depends(get_db)
):
    """
    Génère le document final dans le format demandé.
    - mode="texte"  : utilise le texte corrigé (ou brut si non corrigé)
    - mode="images" : insère les images directement dans le PDF (sans texte)

    Formats supportés : pdf, docx, xlsx, pptx, txt, md
    """
    if format_sortie.lower() not in FORMATS_SUPPORTES:
        raise HTTPException(400, f"Format '{format_sortie}' non supporté. "
                                 f"Choisissez parmi : {', '.join(FORMATS_SUPPORTES)}")

    doc = db.query(DocumentOCR).filter(DocumentOCR.id_ocr == id_ocr).first()
    if not doc:
        raise HTTPException(404, "Document OCR introuvable")

    # Texte à utiliser : corrigé en priorité, sinon brut
    texte = doc.texte_corrige or doc.texte_extrait or ""

    # Reconstruire la liste des images depuis la BDD
    import json as _json
    try:
        chemins_src = _json.loads(doc.chemin_source.replace("'", '"'))
        images = [Path(c) for c in chemins_src if Path(c).exists()]
    except Exception:
        images = []

    mode_images = (mode == "images") and bool(images)
    if mode_images and format_sortie.lower() != "pdf":
        raise HTTPException(400, "Le mode 'images directes' n'est disponible qu'en format PDF.")

    # Générer
    try:
        # Extraction layout avec figures pour le PDF
        if format_sortie.lower() == "pdf" and not mode_images and images:
            blocs_layout_gen = []
            for img_ch in images:
                lr = extraire_regions(img_ch)
                blocs_layout_gen.extend(lr["blocs"])
        else:
            blocs_layout_gen = None
        chemin_dest = dispatcher_generation(format_sortie, texte, titre, images,
                                             mode_images, blocs_layout=blocs_layout_gen)
    except HTTPException:
        raise
    except Exception as e:
        doc.statut = "ECHEC"; doc.message_erreur = str(e)
        db.commit()
        raise HTTPException(500, f"Erreur de génération : {e}")

    # Upload vers MinIO et stocker la clé (pas le chemin local)
    taille = chemin_dest.stat().st_size if chemin_dest.exists() else 0
    objet_key = uploader_fichier_genere(chemin_dest)  # supprime le fichier local
    doc.chemin_pdf     = objet_key
    doc.format_sortie  = format_sortie.lower()
    doc.taille_fichier = taille
    doc.statut         = "GENERE"
    doc.date_traitement = datetime.now()
    db.commit()

    return {
        "id_ocr":        id_ocr,
        "statut":        "GENERE",
        "format_sortie": format_sortie,
        "taille":        taille,
        "message":       f"Document {format_sortie.upper()} généré avec succès.",
        "telecharger_url": f"/ocr/documents/{id_ocr}/telecharger"
    }


# ── 4. TÉLÉCHARGER LE DOCUMENT GÉNÉRÉ ─────────────────────────────────────────
@app.get("/ocr/documents/{id_ocr}/telecharger",
         summary="Télécharger le document généré")
async def telecharger_document(
    id_ocr: int,
    db: Session = Depends(get_db)
):
    doc = db.query(DocumentOCR).filter(DocumentOCR.id_ocr == id_ocr).first()
    if not doc or not doc.chemin_pdf:
        raise HTTPException(404, "Document introuvable ou non encore généré")

    mime_map = {
        "pdf":  "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "txt":  "text/plain; charset=utf-8",
        "md":   "text/markdown; charset=utf-8",
    }
    fmt  = (doc.format_sortie or "pdf").lower()
    mime = mime_map.get(fmt, "application/octet-stream")
    nom  = f"UBH_{doc.nom_fichier_orig or 'document'}_{id_ocr}.{fmt}"

    # Téléchargement depuis MinIO (clé stockée dans doc.chemin_pdf)
    try:
        contenu, _ = telecharger_fichier_genere(doc.chemin_pdf)
    except FileNotFoundError:
        raise HTTPException(404, "Fichier introuvable dans le stockage MinIO")

    return StreamingResponse(
        io.BytesIO(contenu),
        media_type=mime,
        headers={"Content-Disposition": f'attachment; filename="{nom}"'},
    )


# ── 5. WORKFLOW TOUT-EN-UN : SCAN → GÉNÉRER (sans correction) ─────────────────
@app.post("/ocr/scan", summary="Scanner UNE image et générer directement le PDF")
async def scan_simple(
    file:          UploadFile         = File(...),
    langue:        str                = Form("fra"),
    id_doc_ged:    Optional[int]      = Form(None),
    authorization: Optional[str]      = Header(None),
    db:            Session            = Depends(get_db)
):
    """
    Endpoint de compatibilité — scan simple (1 image → PDF immédiat).
    Utilisé par le délégué pour les scans rapides sans workflow de correction.
    """
    if file.content_type not in IMAGE_MIMES:
        raise HTTPException(400, "Format image non supporté")

    user       = await verifier_token(authorization)
    session_id = uuid.uuid4().hex[:12]
    chemin     = sauvegarder_image(file, 0, session_id)
    texte      = ocr_image(chemin, langue)

    doc = DocumentOCR(
        nom_fichier_orig = file.filename,
        chemin_source    = str([str(chemin)]),
        texte_extrait    = texte,
        nb_pages         = 1, nb_images = 1,
        statut           = "RETRANSCRIT",
        langue           = langue,
        id_doc_ged       = id_doc_ged,
        date_traitement  = datetime.now(),
        id_utilisateur   = _uuid_ou_none(user.get("id")) if user else None,
        role_soumetteur  = user.get("role") if user else None,
    )
    db.add(doc); db.commit(); db.refresh(doc)

    # Générer PDF immédiatement
    try:
        chemin_pdf = generer_pdf(texte, file.filename or "Document OCR")
        taille_pdf = chemin_pdf.stat().st_size
        objet_key  = uploader_fichier_genere(chemin_pdf)  # upload MinIO + nettoyage local
        doc.chemin_pdf     = objet_key
        doc.format_sortie  = "pdf"
        doc.taille_fichier = taille_pdf
        doc.statut         = "GENERE"
        db.commit()
    except Exception as e:
        doc.statut = "ECHEC"; doc.message_erreur = str(e)
        db.commit()
        raise HTTPException(500, f"Génération PDF échouée : {e}")

    return {
        "id_ocr":        doc.id_ocr,
        "statut":        "GENERE",
        "texte_extrait": texte,
        "nb_pages":      1,
        "taille_fichier": doc.taille_fichier,
        "telecharger_url": f"/ocr/documents/{doc.id_ocr}/telecharger"
    }


# ── 6. SOUMETTRE VERS GED (avec classement complet) ───────────────────────────
@app.post("/ocr/soumettre-ged",
          summary="Soumettre un document OCR vers le GED avec métadonnées de classement")
async def soumettre_vers_ged(
    background_tasks: BackgroundTasks,
    id_ocr:        int           = Form(...),
    titre:         str           = Form(...),
    type_ressource: str          = Form(...),   # COURS | TD | EXAMEN | ARCHIVE
    id_filiere:    Optional[str] = Form(None),   # str pour accepter "" sans 422
    id_module:     Optional[str] = Form(None),    # str pour accepter "" sans 422
    description:   str           = Form(""),
    depot_direct:  bool          = Form(False), # True = pas besoin de validation
    authorization: Optional[str] = Header(None),
    db:            Session       = Depends(get_db)
):
    """
    Soumet un document OCR déjà généré vers le GED avec toutes ses métadonnées.
    - Si depot_direct=True et rôle autorisé : statut GED = VALIDE immédiatement
    - Sinon : statut = EN_ATTENTE (validation sous-admin / chef-dept)
    """
    user = await verifier_token(authorization)
    if not user:
        raise HTTPException(401, "Authentification requise")

    role = user.get("role", "")
    id_user = user.get("id")

    # Vérification droits de soumission
    if role not in ROLES_SOUMISSION:
        raise HTTPException(403, f"Rôle '{role}' non autorisé à soumettre des documents")

    # Dépôt direct : uniquement pour les rôles privilégiés
    if depot_direct and role not in ROLES_DEPOT_DIRECT:
        depot_direct = False  # Forcer validation même si demandé

    # Convertir id_filiere / id_module (reçus en string depuis FormData, "" → None)
    id_filiere_int: Optional[int] = int(id_filiere) if id_filiere and str(id_filiere).strip().isdigit() else None
    id_module_int:  Optional[int] = int(id_module)  if id_module  and str(id_module).strip().isdigit()  else None
    id_filiere = id_filiere_int   # réassigner pour le reste de la fonction
    id_module  = id_module_int

    doc_ocr = db.query(DocumentOCR).filter(DocumentOCR.id_ocr == id_ocr).first()
    if not doc_ocr:
        raise HTTPException(404, "Document OCR introuvable")
    if doc_ocr.statut not in ("GENERE", "CORRIGE", "RETRANSCRIT"):
        raise HTTPException(400, f"Document en statut '{doc_ocr.statut}' — impossible de soumettre")

    chemin_pdf = Path(doc_ocr.chemin_pdf or "")
    # chemin_pdf est maintenant une CLÉ MinIO (ex: "doc_abc123.pdf")
    # On télécharge le contenu depuis MinIO pour l'envoyer à service-ged
    try:
        contenu_pdf, mime_pdf = telecharger_fichier_genere(doc_ocr.chemin_pdf)
    except FileNotFoundError:
        raise HTTPException(404, "Fichier PDF introuvable dans MinIO — générez d'abord le document")

    nom_pdf = chemin_pdf.name or "document.pdf"

    # Préparer le statut GED
    statut_ged = "VALIDE" if depot_direct else "EN_ATTENTE"
    est_valide = depot_direct

    # Uploader vers service-ged
    token = (authorization or "").replace("Bearer ", "")
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            r = await client.post(
                    f"{API_GED}/ged/documents/upload",
                    files={"file": (nom_pdf, io.BytesIO(contenu_pdf), mime_pdf)},
                    data={
                        "titre":         titre,
                        "type_ressource": type_ressource,
                        "id_filiere":    str(id_filiere) if id_filiere else "",
                        "id_module":     str(id_module)  if id_module  else "",
                        "description":   description,
                        "statut":        statut_ged,
                        "est_valide":    str(est_valide).lower(),
                        "texte_ocr":     doc_ocr.texte_corrige or doc_ocr.texte_extrait or "",
                        "id_soumis_par": str(id_user) if id_user else "",
                    },
                    headers={"Authorization": f"Bearer {token}"}
                )
        if r.status_code not in (200, 201):
            raise HTTPException(r.status_code,
                                f"Erreur GED : {r.text[:200]}")
        doc_ged = r.json()
        id_doc_ged = doc_ged.get("id_document") or doc_ged.get("id_doc")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(502, f"Service GED inaccessible : {e}")

    # Mettre à jour le document OCR
    doc_ocr.id_doc_ged   = id_doc_ged
    doc_ocr.statut       = "SOUMIS" if not depot_direct else "VALIDE"
    doc_ocr.depot_direct = depot_direct
    db.commit()

    # Notification en arrière-plan
    if id_user and token:
        msg = (
            f"Votre document '{titre}' a été déposé directement et est maintenant disponible."
            if depot_direct else
            f"Votre document '{titre}' a été soumis avec succès et est en attente de validation."
        )
        background_tasks.add_task(
            notifier_utilisateur, id_user,
            "Soumission confirmée" if not depot_direct else "Document publié",
            msg, token
        )

    return {
        "id_ocr":      id_ocr,
        "id_doc_ged":  id_doc_ged,
        "statut_ocr":  doc_ocr.statut,
        "statut_ged":  statut_ged,
        "depot_direct": depot_direct,
        "message": (
            f"Document déposé directement — disponible immédiatement."
            if depot_direct else
            f"Document soumis — en attente de validation par le sous-administrateur."
        )
    }


# ── 7. VALIDER / REFUSER UN DOCUMENT (sous-admin / chef-dept) ─────────────────
@app.put("/ocr/documents/{id_ocr}/valider",
         summary="Valider ou refuser un document OCR soumis (sous-admin / chef-dept)")
async def valider_document_ocr(
    id_ocr:        int,
    background_tasks: BackgroundTasks,
    statut:        str           = Form(...),    # VALIDE | REFUSE
    motif_refus:   str           = Form(""),
    titre:         Optional[str] = Form(None),
    type_ressource: Optional[str]= Form(None),
    id_filiere:    Optional[int] = Form(None),
    id_module:     Optional[int] = Form(None),
    authorization: Optional[str] = Header(None),
    db:            Session       = Depends(get_db)
):
    """
    Valide ou refuse un document soumis.
    - Si VALIDE : met à jour le GED (statut + classement)
    - Si REFUSE  : notifie le soumetteur avec le motif
    """
    user = await verifier_token(authorization)
    if not user:
        raise HTTPException(401, "Authentification requise")

    role = user.get("role", "")
    if role not in {"ADMIN", "SOUS_ADMIN", "CHEF_DEPARTEMENT", "PRESIDENT"}:
        raise HTTPException(403, "Action réservée aux sous-administrateurs et chefs de département")

    if statut not in ("VALIDE", "REFUSE"):
        raise HTTPException(400, "Statut invalide — choisir VALIDE ou REFUSE")
    if statut == "REFUSE" and not motif_refus.strip():
        raise HTTPException(400, "Motif de refus obligatoire")

    doc_ocr = db.query(DocumentOCR).filter(DocumentOCR.id_ocr == id_ocr).first()
    if not doc_ocr:
        raise HTTPException(404, "Document OCR introuvable")

    token = (authorization or "").replace("Bearer ", "")

    # Mettre à jour le GED si document lié
    if doc_ocr.id_doc_ged:
        try:
            async with httpx.AsyncClient(timeout=10.0) as client:
                await client.put(
                    f"{API_GED}/ged/documents/{doc_ocr.id_doc_ged}/valider",
                    json={
                        "statut":        statut,
                        "motif_refus":   motif_refus if statut == "REFUSE" else None,
                        "titre":         titre,
                        "type_ressource": type_ressource,
                        "id_filiere":    id_filiere,
                        "id_module":     id_module,
                    },
                    headers={"Authorization": f"Bearer {token}"}
                )
        except Exception as e:
            log.warning(f"Mise à jour GED échouée pour doc {doc_ocr.id_doc_ged} : {e}")

    # Mettre à jour le document OCR
    doc_ocr.statut = statut
    db.commit()

    # Notifier le soumetteur
    if doc_ocr.id_utilisateur and token:
        if statut == "VALIDE":
            titre_notif = "Document validé ✓"
            message = (
                f"Votre document '{doc_ocr.nom_fichier_orig}' a été validé "
                f"et est maintenant accessible à tous les étudiants de votre filière."
            )
        else:
            titre_notif = "Document refusé"
            message = (
                f"Votre document '{doc_ocr.nom_fichier_orig}' a été refusé. "
                f"Motif : {motif_refus}"
            )
        background_tasks.add_task(
            notifier_utilisateur,
            doc_ocr.id_utilisateur, titre_notif, message, token
        )

    return {
        "id_ocr":    id_ocr,
        "statut":    statut,
        "message":   f"Document {statut.lower()} avec succès."
    }



# ── /ocr/retranscrire  (appelé par ocr_studio.html) — v3.0 PaddleOCR uniquement ──
@app.post("/ocr/retranscrire",
          summary="Retranscrire N images — PaddleOCR + détection layout OpenCV")
async def retranscrire(
    files:         List[UploadFile] = File(...),
    langue:        str              = Form("fra"),
    ordre:         str              = Form(""),
    authorization: Optional[str]   = Header(None),
    db:            Session          = Depends(get_db)
):
    user       = await verifier_token(authorization)
    session_id = uuid.uuid4().hex[:12]
    chemins    = []

    for i, f in enumerate(files):
        if f.content_type not in IMAGE_MIMES:
            raise HTTPException(400, f"{f.filename} : format non supporte")
        chemin = sauvegarder_image(f, i, session_id)
        chemins.append(chemin)

    if ordre.strip():
        try:
            indices = [int(x) for x in ordre.split(",")]
            if len(indices) == len(chemins):
                chemins = [chemins[i] for i in indices]
        except Exception:
            pass

    # ── OCR sur chaque image : Gemini Vision en priorité, PaddleOCR en fallback ─
    texte_global       = ""
    resultats          = []
    moteurs_utilises   = []
    has_figures_global = False

    for i, chemin in enumerate(chemins):
        texte_img   = ""
        has_figures = False
        moteur_img  = "paddleocr"

        # ── PaddleOCR + détection layout OpenCV ──────────────────────────────
        try:
            layout_result = extraire_regions(chemin)
            blocs         = layout_result["blocs"]
            has_figures   = any(b["type"] == "figure" for b in blocs)
            texte_paddle  = "\n".join(b["texte"] for b in blocs if b["type"] == "texte")
            texte_img     = _nettoyer_texte_ocr(texte_paddle)
            moteur_img    = "paddleocr+layout" if has_figures else "paddleocr"
            log.info(f"retranscrire — image {i+1} — PaddleOCR OK ({len(texte_img)} chars)")
        except HTTPException as e:
            log.error(f"retranscrire — image {i+1} — PaddleOCR erreur {e.status_code}: {e.detail}")
            texte_img  = ""
            moteur_img = "echec"

        moteurs_utilises.append(moteur_img)
        has_figures_global = has_figures_global or has_figures
        score = _score_confiance(texte_img, moteur_img, has_figures)
        resultats.append({
            "index":         i,
            "nom_fichier":   chemin.name,
            "texte_extrait": texte_img,
            "nb_mots":       len(texte_img.split()),
            "has_figures":   has_figures,
            "moteur":        moteur_img,
            "score_confiance": score,
        })
        texte_global += texte_img + "\n\n"

    texte_global = texte_global.strip()

    # ── Résumé moteur principal ────────────────────────────────────────────────
    moteur_principal = (
        "paddleocr+layout" if all(m == "paddleocr+layout" for m in moteurs_utilises)
        else (moteurs_utilises[0] if moteurs_utilises else "paddleocr")
    )

    score_global = (
        sum(r["score_confiance"] for r in resultats) // len(resultats)
        if resultats else 0
    )

    doc = DocumentOCR(
        nom_fichier_orig = files[0].filename if files else "scan",
        chemin_source    = str([str(c) for c in chemins]),
        texte_extrait    = texte_global,
        texte_corrige    = texte_global,
        nb_pages         = len(files),
        nb_images        = len(files),
        statut           = "RETRANSCRIT",
        langue           = langue,
        date_traitement  = datetime.now(),
        id_utilisateur   = _uuid_ou_none(user.get("id")) if user else None,
        role_soumetteur  = user.get("role") if user else None,
    )
    db.add(doc); db.commit(); db.refresh(doc)

    return {
        "id_ocr":              doc.id_ocr,
        "nb_images":           len(files),
        "nb_mots":             len(texte_global.split()),
        "langue":              langue,
        "statut":              doc.statut,
        "moteur_ocr":          moteur_principal,
        "has_figures":         has_figures_global,
        "score_confiance":     score_global,
        "resultats_par_image": resultats,
        "texte_consolide":     texte_global,
        "texte_extrait":       texte_global,
        "texte_corrige":       texte_global,
        "message": (
            f"{len(files)} image(s) retranscrite(s) par {moteur_principal}"
            + (" (layout/figures détectés)" if has_figures_global else "")
            + f" — confiance : {score_global}%."
        )
    }


# ── /ocr/generer-document  (appelé par ocr_studio.html) — v3.0 PaddleOCR uniquement ──
@app.post("/ocr/generer-document",
          summary="Generer document final et retourner le fichier directement (blob)")
async def generer_document_direct(
    files:         List[UploadFile] = File(...),
    titre:         str              = Form("Document OCR"),
    format_sortie: str              = Form("pdf"),
    texte_corrige: str              = Form(""),
    mode_images:   str              = Form("false"),
    langue:        str              = Form("fra"),
    authorization: Optional[str]   = Header(None),
    db:            Session          = Depends(get_db)
):
    fmt = format_sortie.lower().strip()
    if fmt not in FORMATS_SUPPORTES:
        raise HTTPException(400, f"Format '{fmt}' non supporte. Choisissez : {FORMATS_SUPPORTES}")

    user       = await verifier_token(authorization)
    session_id = uuid.uuid4().hex[:12]
    chemins    = []

    for i, f in enumerate(files):
        if f.content_type not in IMAGE_MIMES:
            raise HTTPException(400, f"{f.filename} : format image non supporte")
        chemin = sauvegarder_image(f, i, session_id)
        chemins.append(chemin)

    use_images = (mode_images.lower() == "true")
    if use_images and fmt != "pdf":
        raise HTTPException(400, "Mode images directes uniquement disponible en PDF")

    # ── Déterminer le texte à utiliser ────────────────────────────────────────
    if texte_corrige.strip():
        # Texte corrigé fourni par le frontend → priorité absolue
        texte = texte_corrige.strip()
    elif use_images:
        texte = ""
    else:
        # Aucun texte fourni → OCR PaddleOCR sur chaque image
        textes_par_image = []
        for chemin in chemins:
            try:
                layout_result = extraire_regions(chemin)
                blocs         = layout_result["blocs"]
                texte_paddle  = "\n".join(b["texte"] for b in blocs if b["type"] == "texte")
                textes_par_image.append(_nettoyer_texte_ocr(texte_paddle))
                log.info(f"generer-document — PaddleOCR OK ({len(texte_paddle)} chars)")
            except HTTPException:
                raise
            except Exception as e:
                log.error(f"generer-document — OCR échoué : {e}")
                textes_par_image.append("")

        texte = "\n\n".join(textes_par_image).strip()

    try:
        chemin_dest = dispatcher_generation(
            fmt, texte, titre,
            images=chemins if use_images else None,
            mode_images=use_images
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Erreur generation : {e}")

    taille    = chemin_dest.stat().st_size if chemin_dest.exists() else 0
    objet_key = uploader_fichier_genere(chemin_dest)  # upload MinIO + nettoyage local

    doc = DocumentOCR(
        nom_fichier_orig = files[0].filename if files else "document",
        chemin_source    = str([str(c) for c in chemins]),
        chemin_pdf       = objet_key,
        texte_extrait    = texte,
        texte_corrige    = texte_corrige.strip() or None,
        nb_pages         = len(files),
        nb_images        = len(files),
        format_sortie    = fmt,
        taille_fichier   = taille,
        statut           = "GENERE",
        langue           = langue,
        date_traitement  = datetime.now(),
        id_utilisateur   = _uuid_ou_none(user.get("id")) if user else None,
        role_soumetteur  = user.get("role") if user else None,
    )
    db.add(doc); db.commit(); db.refresh(doc)

    mime_map = {
        "pdf":  "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "txt":  "text/plain; charset=utf-8",
        "md":   "text/markdown; charset=utf-8",
    }
    mime     = mime_map.get(fmt, "application/octet-stream")
    filename = (titre.replace(" ", "_")[:40] or "document") + "." + fmt

    # Stream depuis MinIO
    try:
        contenu, _ = telecharger_fichier_genere(objet_key)
    except FileNotFoundError:
        raise HTTPException(500, "Fichier généré introuvable dans MinIO")

    return StreamingResponse(
        io.BytesIO(contenu),
        media_type=mime,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )






# ══════════════════════════════════════════════════════════════════════════════
# NOUVEAUX ENDPOINTS — SOUMISSION ET GÉNÉRATION DIRECTE (SANS OCR PRÉALABLE)
# Permettent d'envoyer n'importe quel fichier (image, PDF, DOC, etc.)
# directement vers la GED ou de générer un document sans passer par l'OCR.
# ══════════════════════════════════════════════════════════════════════════════

# Tous les types MIME acceptés en soumission directe
DIRECT_MIMES = {
    *IMAGE_MIMES,
    "application/pdf",
    "application/msword",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.ms-powerpoint",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-excel",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "text/plain",
    "text/markdown",
}


def _sauvegarder_fichier_direct(upload: UploadFile, idx: int, session_id: str) -> Path:
    """Sauvegarde un fichier quelconque (pas seulement image) dans le répertoire upload."""
    ext  = Path(upload.filename or f"fichier_{idx}").suffix or ".bin"
    dest = UPLOAD_DIR / f"{session_id}_direct_{idx}{ext}"
    import shutil
    upload.file.seek(0)
    with open(str(dest), "wb") as fh:
        shutil.copyfileobj(upload.file, fh)
    return dest


@app.post(
    "/ocr/generer-direct",
    summary="Générer un document depuis fichiers bruts SANS OCR préalable",
)
async def generer_direct(
    fichiers:      List[UploadFile] = File(...),
    titre:         str              = Form("Document"),
    format_sortie: str              = Form("pdf"),
    texte_corrige: str              = Form(""),
    authorization: Optional[str]   = Header(None),
    db:            Session          = Depends(get_db),
):
    """
    Génère un document (PDF, DOCX, XLSX, PPTX, TXT, MD) à partir de fichiers bruts.
    Si les fichiers sont des images ET qu'aucun texte n'est fourni, effectue l'OCR.
    Si les fichiers sont des non-images (PDF, DOC…), les joint directement au PDF.
    Retourne le fichier en téléchargement direct (Content-Disposition: attachment).
    """
    fmt = format_sortie.lower().strip()
    if fmt not in FORMATS_SUPPORTES:
        raise HTTPException(400, f"Format '{fmt}' non supporte. Choisissez : {FORMATS_SUPPORTES}")

    user       = await verifier_token(authorization)
    session_id = uuid.uuid4().hex[:12]
    chemins    = []
    noms       = []

    for i, f in enumerate(fichiers):
        chemin = _sauvegarder_fichier_direct(f, i, session_id)
        chemins.append(chemin)
        noms.append(f.filename or f"fichier_{i}")

    # Séparer images et autres fichiers
    chemins_images = [
        c for c, f in zip(chemins, fichiers)
        if (f.content_type or "") in IMAGE_MIMES
        or Path(f.filename or "").suffix.lower() in {".jpg", ".jpeg", ".png", ".tiff", ".bmp", ".webp"}
    ]

    # ── Génération DIRECTE : jamais d'OCR automatique ─────────────────────────
    # Si l'utilisateur n'a pas retranscrit, on insère les images telles quelles
    # dans le document (mode_images=True pour PDF, sinon texte minimal).
    texte = texte_corrige.strip()  # vide si pas de retranscription préalable

    if not texte and not chemins_images:
        # Fichiers non-image sans texte : titre comme contenu minimal
        texte = titre

    # Pour les images sans texte : mode_images activé (images dans le PDF dans l'ordre)
    mode_images_direct = bool(chemins_images) and not texte and fmt == "pdf"

    # Générer le document
    try:
        chemin_dest = dispatcher_generation(
            fmt, texte, titre,
            images=chemins_images if chemins_images else None,
            mode_images=mode_images_direct,
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Erreur generation : {e}")

    taille    = chemin_dest.stat().st_size if chemin_dest.exists() else 0
    objet_key = uploader_fichier_genere(chemin_dest)

    # Créer un enregistrement DocumentOCR pour le suivi
    doc = DocumentOCR(
        nom_fichier_orig = noms[0] if noms else "document",
        chemin_source    = str([str(c) for c in chemins]),
        chemin_pdf       = objet_key,
        texte_extrait    = texte,
        texte_corrige    = texte_corrige.strip() or None,
        nb_pages         = len(fichiers),
        nb_images        = len(chemins_images),
        format_sortie    = fmt,
        taille_fichier   = taille,
        statut           = "GENERE",
        langue           = "fra",
        date_traitement  = datetime.now(),
        id_utilisateur   = _uuid_ou_none(user.get("id")) if user else None,
        role_soumetteur  = user.get("role") if user else None,
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)

    # Renvoyer le fichier en téléchargement direct
    mime_map = {
        "pdf":  "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "txt":  "text/plain; charset=utf-8",
        "md":   "text/markdown; charset=utf-8",
    }
    mime     = mime_map.get(fmt, "application/octet-stream")
    filename = (titre.replace(" ", "_")[:40] or "document") + "." + fmt

    try:
        contenu, _ = telecharger_fichier_genere(objet_key)
    except FileNotFoundError:
        raise HTTPException(500, "Fichier genere introuvable apres upload MinIO")

    return StreamingResponse(
        io.BytesIO(contenu),
        media_type=mime,
        headers={
            "Content-Disposition": f'attachment; filename="{filename}"',
            "X-Id-Ocr": str(doc.id_ocr),
        },
    )


@app.post(
    "/ocr/soumettre-direct",
    summary="Soumettre des fichiers bruts vers la GED SANS OCR préalable",
)
async def soumettre_direct(
    background_tasks: BackgroundTasks,
    fichiers:      List[UploadFile] = File(...),
    titre:         str              = Form(...),
    type_ressource: str             = Form(...),
    id_filiere:    Optional[int]    = Form(None),
    id_module:     Optional[int]    = Form(None),
    description:   str              = Form(""),
    format_sortie: str              = Form("pdf"),
    texte_corrige: str              = Form(""),
    depot_direct:  bool             = Form(False),
    authorization: Optional[str]    = Header(None),
    db:            Session          = Depends(get_db),
):
    """
    Soumet N fichiers (images, PDF, DOC, PPT, XLS, TXT…) directement vers la GED
    sans OCR préalable.
    Workflow :
      1. Sauvegarde les fichiers localement
      2. Génère un PDF composite (ou utilise le premier PDF brut si un seul PDF fourni)
      3. Upload vers MinIO
      4. Crée l'entrée DocumentOCR en BDD (statut GENERE)
      5. Envoie vers service-ged pour validation
    """
    user = await verifier_token(authorization)
    if not user:
        raise HTTPException(401, "Authentification requise")

    role    = user.get("role", "")
    id_user = user.get("id")

    if role not in ROLES_SOUMISSION:
        raise HTTPException(403, f"Role '{role}' non autorise a soumettre des documents")

    if depot_direct and role not in ROLES_DEPOT_DIRECT:
        depot_direct = False

    session_id = uuid.uuid4().hex[:12]
    chemins    = []
    noms       = []

    for i, f in enumerate(fichiers):
        chemin = _sauvegarder_fichier_direct(f, i, session_id)
        chemins.append(chemin)
        noms.append(f.filename or f"fichier_{i}")

    # Séparer images des autres
    chemins_images = [
        c for c, f in zip(chemins, fichiers)
        if (f.content_type or "") in IMAGE_MIMES
        or Path(f.filename or "").suffix.lower() in {".jpg", ".jpeg", ".png", ".tiff", ".bmp", ".webp"}
    ]

    # Cas particulier : un seul PDF fourni → on l'utilise directement
    chemins_pdf_bruts = [
        c for c, f in zip(chemins, fichiers)
        if (f.content_type or "") == "application/pdf"
        or Path(f.filename or "").suffix.lower() == ".pdf"
    ]

    # ── Soumission DIRECTE : jamais d'OCR automatique ─────────────────────────
    # Si l'utilisateur n'a pas retranscrit, on soumet les images telles quelles.
    texte = texte_corrige.strip()  # vide si pas de retranscription préalable

    if not texte and not chemins_images:
        texte = titre

    # Choisir le fichier à envoyer à la GED
    fmt = format_sortie.lower().strip() if format_sortie.lower().strip() in FORMATS_SUPPORTES else "pdf"

    # Mode images : PDF contenant les images dans l'ordre, sans OCR
    mode_images_direct = bool(chemins_images) and not texte and fmt == "pdf"

    if len(chemins_pdf_bruts) == 1 and len(fichiers) == 1:
        # Un seul PDF brut → l'envoyer tel quel sans re-génération
        chemin_dest = chemins_pdf_bruts[0]
        taille      = chemin_dest.stat().st_size
        objet_key   = uploader_fichier_genere(chemin_dest)
        mime_envoi  = "application/pdf"
        nom_envoi   = chemin_dest.name
        fmt         = "pdf"
    else:
        # Générer un PDF composite (images dans l'ordre, sans OCR)
        try:
            chemin_dest = dispatcher_generation(
                fmt, texte, titre,
                images=chemins_images if chemins_images else None,
                mode_images=mode_images_direct,
            )
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(500, f"Erreur generation document : {e}")

        taille     = chemin_dest.stat().st_size if chemin_dest.exists() else 0
        objet_key  = uploader_fichier_genere(chemin_dest)
        mime_map   = {
            "pdf":  "application/pdf",
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "txt":  "text/plain; charset=utf-8",
            "md":   "text/markdown; charset=utf-8",
        }
        mime_envoi = mime_map.get(fmt, "application/pdf")
        nom_envoi  = (titre.replace(" ", "_")[:40] or "document") + "." + fmt

    # Créer l'entrée BDD
    doc = DocumentOCR(
        nom_fichier_orig = noms[0] if noms else "document",
        chemin_source    = str([str(c) for c in chemins]),
        chemin_pdf       = objet_key,
        texte_extrait    = texte,
        texte_corrige    = texte_corrige.strip() or None,
        nb_pages         = len(fichiers),
        nb_images        = len(chemins_images),
        format_sortie    = fmt,
        taille_fichier   = taille,
        statut           = "GENERE",
        langue           = "fra",
        date_traitement  = datetime.now(),
        id_utilisateur   = _uuid_ou_none(id_user),
        role_soumetteur  = role,
        depot_direct     = depot_direct,
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)
    id_ocr = doc.id_ocr

    # Télécharger le contenu pour l'envoyer au service-ged
    try:
        contenu_fichier, _ = telecharger_fichier_genere(objet_key)
    except FileNotFoundError:
        raise HTTPException(500, "Fichier introuvable apres upload MinIO")

    statut_ged = "VALIDE" if depot_direct else "EN_ATTENTE"
    est_valide = depot_direct
    token      = (authorization or "").replace("Bearer ", "")

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            r = await client.post(
                f"{API_GED}/ged/documents/upload",
                files={"file": (nom_envoi, io.BytesIO(contenu_fichier), mime_envoi)},
                data={
                    "titre":          titre,
                    "type_ressource": type_ressource,
                    "id_filiere":     str(id_filiere) if id_filiere else "",
                    "id_module":      str(id_module)  if id_module  else "",
                    "description":    description,
                    "statut":         statut_ged,
                    "est_valide":     str(est_valide).lower(),
                    "texte_ocr":      texte,
                    "id_soumis_par":  str(id_user) if id_user else "",
                },
                headers={"Authorization": f"Bearer {token}"},
            )
        if r.status_code not in (200, 201):
            raise HTTPException(r.status_code, f"Erreur GED : {r.text[:200]}")
        doc_ged    = r.json()
        id_doc_ged = doc_ged.get("id_document") or doc_ged.get("id_doc")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(502, f"Service GED inaccessible : {e}")

    # Mettre à jour le statut OCR
    doc.id_doc_ged = id_doc_ged
    doc.statut     = "SOUMIS" if not depot_direct else "VALIDE"
    db.commit()

    # Notification en arrière-plan
    if id_user and token:
        msg = (
            f"Votre document '{titre}' a ete depose directement et est maintenant disponible."
            if depot_direct else
            f"Votre document '{titre}' a ete soumis avec succes et est en attente de validation."
        )
        background_tasks.add_task(
            notifier_utilisateur,
            id_user,
            "Soumission confirmee" if not depot_direct else "Document publie",
            msg,
            token,
        )

    return {
        "id_ocr":       id_ocr,
        "id_doc_ged":   id_doc_ged,
        "statut_ocr":   doc.statut,
        "statut_ged":   statut_ged,
        "depot_direct": depot_direct,
        "nb_fichiers":  len(fichiers),
        "message": (
            "Document depose directement — disponible immediatement."
            if depot_direct else
            "Document soumis — en attente de validation."
        ),
    }


# ── 8. SCAN SIMPLE + GED (rétro-compatibilité) ────────────────────────────────
@app.post("/ocr/scan-et-ged",
          summary="Scanner une image et soumettre automatiquement vers le GED")
async def scan_et_ged(
    background_tasks: BackgroundTasks,
    file:          UploadFile         = File(...),
    langue:        str                = Form("fra"),
    id_doc_ged:    Optional[int]      = Form(None),
    authorization: Optional[str]      = Header(None),
    db:            Session            = Depends(get_db)
):
    """Rétro-compatibilité : scan 1 image + lien vers un doc GED existant."""
    if file.content_type not in IMAGE_MIMES:
        raise HTTPException(400, "Format image non supporté")

    user       = await verifier_token(authorization)
    session_id = uuid.uuid4().hex[:12]
    chemin     = sauvegarder_image(file, 0, session_id)
    texte      = ocr_image(chemin, langue)

    # Générer PDF + upload MinIO
    chemin_pdf  = generer_pdf(texte, file.filename or "Document OCR")
    taille_pdf  = chemin_pdf.stat().st_size
    objet_key   = uploader_fichier_genere(chemin_pdf)  # upload MinIO + nettoyage local

    doc = DocumentOCR(
        nom_fichier_orig = file.filename,
        chemin_source    = str([str(chemin)]),
        chemin_pdf       = objet_key,
        texte_extrait    = texte,
        nb_pages         = 1, nb_images = 1,
        format_sortie    = "pdf",
        taille_fichier   = taille_pdf,
        statut           = "GENERE",
        langue           = langue,
        id_doc_ged       = id_doc_ged,
        date_traitement  = datetime.now(),
        id_utilisateur   = _uuid_ou_none(user.get("id")) if user else None,
        role_soumetteur  = user.get("role") if user else None,
    )
    db.add(doc); db.commit(); db.refresh(doc)

    # Injecter le texte dans le doc GED si id_doc_ged fourni
    if id_doc_ged:
        token = (authorization or "").replace("Bearer ", "")
        try:
            async with httpx.AsyncClient(timeout=10.0) as client:
                await client.patch(
                    f"{API_GED}/ged/documents/{id_doc_ged}/texte-ocr",
                    json={"texte_ocr": texte},
                    headers={"Authorization": f"Bearer {token}"}
                )
        except Exception:
            pass

    return {
        "id_ocr":        doc.id_ocr,
        "statut":        "GENERE",
        "texte_extrait": texte,
        "nb_pages":      1,
        "taille_fichier": doc.taille_fichier,
        "message":       "Scan et synchronisation GED effectués."
    }


# ── 9. LISTE DOCUMENTS OCR ────────────────────────────────────────────────────
@app.get("/ocr/documents", summary="Liste des documents OCR")
async def liste_documents(
    statut:        Optional[str] = None,
    limite:        int           = 50,
    offset:        int           = 0,
    authorization: Optional[str] = Header(None),
    db:            Session       = Depends(get_db)
):
    q = db.query(DocumentOCR)
    if statut:
        q = q.filter(DocumentOCR.statut == statut)
    total = q.count()
    docs  = q.order_by(DocumentOCR.date_soumission.desc()).offset(offset).limit(limite).all()
    return {
        "total": total,
        "limite": limite,
        "offset": offset,
        "documents": [
            {
                "id_ocr":           d.id_ocr,
                "nom_fichier_orig": d.nom_fichier_orig,
                "format_sortie":    d.format_sortie,
                "nb_pages":         d.nb_pages,
                "nb_images":        d.nb_images,
                "taille_fichier":   d.taille_fichier,
                "statut":           d.statut,
                "langue":           d.langue,
                "depot_direct":     d.depot_direct,
                "id_doc_ged":       d.id_doc_ged,
                "date_soumission":  d.date_soumission.isoformat() if d.date_soumission else None,
                "date_traitement":  d.date_traitement.isoformat()  if d.date_traitement  else None,
                "message_erreur":   d.message_erreur,
                "telecharger_url":  f"/ocr/documents/{d.id_ocr}/telecharger" if d.chemin_pdf else None,
            }
            for d in docs
        ]
    }


# ── 10. DÉTAIL D'UN DOCUMENT OCR ──────────────────────────────────────────────
@app.get("/ocr/documents/{id_ocr}", summary="Détail d'un document OCR")
async def detail_document(id_ocr: int, db: Session = Depends(get_db)):
    doc = db.query(DocumentOCR).filter(DocumentOCR.id_ocr == id_ocr).first()
    if not doc:
        raise HTTPException(404, "Document OCR introuvable")
    return {
        "id_ocr":           doc.id_ocr,
        "nom_fichier_orig": doc.nom_fichier_orig,
        "format_sortie":    doc.format_sortie,
        "texte_extrait":    doc.texte_extrait,
        "texte_corrige":    doc.texte_corrige,
        "nb_pages":         doc.nb_pages,
        "nb_images":        doc.nb_images,
        "taille_fichier":   doc.taille_fichier,
        "statut":           doc.statut,
        "langue":           doc.langue,
        "depot_direct":     doc.depot_direct,
        "id_doc_ged":       doc.id_doc_ged,
        "date_soumission":  doc.date_soumission.isoformat() if doc.date_soumission else None,
        "date_traitement":  doc.date_traitement.isoformat()  if doc.date_traitement  else None,
        "message_erreur":   doc.message_erreur,
        "telecharger_url":  f"/ocr/documents/{doc.id_ocr}/telecharger" if doc.chemin_pdf else None,
    }


# ── 11. STATS OCR ─────────────────────────────────────────────────────────────
@app.get("/ocr/stats", summary="Statistiques du service OCR")
async def stats_ocr(db: Session = Depends(get_db)):
    from sqlalchemy import func as sqlfunc
    total       = db.query(DocumentOCR).count()
    termines    = db.query(DocumentOCR).filter(
                      DocumentOCR.statut.in_(["GENERE","SOUMIS","VALIDE"])).count()
    echecs      = db.query(DocumentOCR).filter(DocumentOCR.statut == "ECHEC").count()
    en_attente  = db.query(DocumentOCR).filter(DocumentOCR.statut == "EN_ATTENTE").count()
    soumis      = db.query(DocumentOCR).filter(DocumentOCR.statut == "SOUMIS").count()
    valides     = db.query(DocumentOCR).filter(DocumentOCR.statut == "VALIDE").count()
    refuses     = db.query(DocumentOCR).filter(DocumentOCR.statut == "REFUSE").count()
    pages_total = db.query(sqlfunc.sum(DocumentOCR.nb_pages)).scalar() or 0
    taille_tot  = db.query(sqlfunc.sum(DocumentOCR.taille_fichier)).scalar() or 0

    formats = {}
    for f in FORMATS_SUPPORTES:
        formats[f] = db.query(DocumentOCR).filter(
            DocumentOCR.format_sortie == f).count()

    return {
        "total_documents":   total,
        "total_termines":    termines,
        "total_echecs":      echecs,
        "total_en_attente":  en_attente,
        "total_soumis":      soumis,
        "total_valides":     valides,
        "total_refuses":     refuses,
        "total_pages":       int(pages_total),
        "taille_totale_mo":  round(taille_tot / 1048576, 2),
        "formats_supportes": FORMATS_SUPPORTES,
        "repartition_formats": formats,
    }



# ── 13. RETRANSCRIPTION EN STREAMING SSE (progression temps réel) ─────────────
@app.post("/ocr/retranscrire-stream",
          summary="Retranscrire N images avec progression SSE — PaddleOCR uniquement")
async def retranscrire_stream(
    files:         List[UploadFile] = File(...),
    langue:        str              = Form("fra"),
    ordre:         str              = Form(""),
    authorization: Optional[str]   = Header(None),
    db:            Session          = Depends(get_db)
):
    """
    Identique à /ocr/retranscrire mais retourne un flux SSE (text/event-stream).
    Chaque image envoie un event 'progress' avec son texte, moteur et score.
    Un event 'done' final envoie l'id_ocr et le texte consolidé.

    Format SSE :
      event: progress
      data: {"image_index":0,"total":3,"nom":"img.jpg","texte":"...","moteur":"paddleocr","score":72,"has_figures":false}

      event: done
      data: {"id_ocr":42,"nb_mots":450,"score_global":70,"texte_consolide":"...","moteur_ocr":"paddleocr"}

      event: error
      data: {"message":"..."}
    """
    import json as _json_sse

    user       = await verifier_token(authorization)
    session_id = uuid.uuid4().hex[:12]
    chemins    = []

    for i, f in enumerate(files):
        if f.content_type not in IMAGE_MIMES:
            async def _err():
                yield f"event: error\ndata: {_json_sse.dumps({'message': f'{f.filename} : format non supporte'})}\n\n"
            return StreamingResponse(_err(), media_type="text/event-stream")
        chemin = sauvegarder_image(f, i, session_id)
        chemins.append(chemin)

    if ordre.strip():
        try:
            indices = [int(x) for x in ordre.split(",")]
            if len(indices) == len(chemins):
                chemins = [chemins[i] for i in indices]
        except Exception:
            pass

    noms_fichiers = [f.filename or f"image_{i}" for i, f in enumerate(files)]
    total = len(chemins)

    async def generateur():
        texte_global       = ""
        resultats          = []
        moteurs_utilises   = []
        has_figures_global = False

        for i, chemin in enumerate(chemins):
            texte_img   = ""
            has_figures = False
            moteur_img  = "paddleocr"

            # Progression : annonce du démarrage de l'image
            yield (
                f"event: progress\n"
                f"data: {_json_sse.dumps({'image_index': i, 'total': total, 'phase': 'debut', 'nom': noms_fichiers[i]})}\n\n"
            )

            # ── PaddleOCR + détection layout OpenCV ─────────────────────────
            yield (
                f"event: progress\n"
                f"data: {_json_sse.dumps({'image_index': i, 'total': total, 'phase': 'paddleocr', 'nom': noms_fichiers[i]})}\n\n"
            )
            try:
                layout_result = extraire_regions(chemin)
                blocs         = layout_result["blocs"]
                has_figures   = any(b["type"] == "figure" for b in blocs)
                texte_paddle  = "\n".join(b["texte"] for b in blocs if b["type"] == "texte")
                texte_img     = _nettoyer_texte_ocr(texte_paddle)
                moteur_img    = "paddleocr+layout" if has_figures else "paddleocr"
            except HTTPException as e:
                log.error(f"SSE image {i+1} — OCR erreur : {e.detail}")
                texte_img  = ""
                moteur_img = "echec"

            score = _score_confiance(texte_img, moteur_img, has_figures)
            moteurs_utilises.append(moteur_img)
            has_figures_global = has_figures_global or has_figures
            resultats.append({
                "index": i, "nom_fichier": noms_fichiers[i],
                "texte_extrait": texte_img, "nb_mots": len(texte_img.split()),
                "has_figures": has_figures, "moteur": moteur_img,
                "score_confiance": score,
            })
            texte_global += texte_img + "\n\n"

            # Résultat de l'image
            yield (
                f"event: progress\n"
                f"data: {_json_sse.dumps({'image_index': i, 'total': total, 'phase': 'termine', 'nom': noms_fichiers[i], 'texte': texte_img, 'moteur': moteur_img, 'score': score, 'has_figures': has_figures})}\n\n"
            )

        texte_global = texte_global.strip()
        moteur_principal = (
            "paddleocr+layout" if all(m == "paddleocr+layout" for m in moteurs_utilises)
            else (moteurs_utilises[0] if moteurs_utilises else "paddleocr")
        )
        score_global = (
            sum(r["score_confiance"] for r in resultats) // len(resultats)
            if resultats else 0
        )

        # Persister en BDD
        try:
            doc = DocumentOCR(
                nom_fichier_orig = noms_fichiers[0] if noms_fichiers else "scan",
                chemin_source    = str([str(c) for c in chemins]),
                texte_extrait    = texte_global,
                texte_corrige    = texte_global,
                nb_pages         = total,
                nb_images        = total,
                statut           = "RETRANSCRIT",
                langue           = langue,
                date_traitement  = datetime.now(),
                id_utilisateur   = _uuid_ou_none(user.get("id")) if user else None,
                role_soumetteur  = user.get("role") if user else None,
            )
            db.add(doc); db.commit(); db.refresh(doc)
            id_ocr_final = doc.id_ocr
        except Exception as e:
            log.error(f"SSE BDD erreur : {e}")
            id_ocr_final = None

        yield (
            f"event: done\n"
            f"data: {_json_sse.dumps({'id_ocr': id_ocr_final, 'nb_images': total, 'nb_mots': len(texte_global.split()), 'score_global': score_global, 'moteur_ocr': moteur_principal, 'has_figures': has_figures_global, 'texte_consolide': texte_global, 'texte_corrige': texte_global, 'resultats_par_image': resultats})}\n\n"
        )

    return StreamingResponse(
        generateur(),
        media_type="text/event-stream",
        headers={
            "Cache-Control":   "no-cache",
            "X-Accel-Buffering": "no",   # Désactiver le buffering nginx si présent
        }
    )


# ── 12. HEALTH CHECK ──────────────────────────────────────────────────────────
@app.get("/health")
@app.get("/ocr/health")   # alias pour le proxy frontend /api/ocr/health
def health():
    from ocr_engine import GROQ_AVAILABLE as _groq
    return {
        "status":  "ok" if PADDLE_AVAILABLE else "degraded",
        "service": "uniburkina-ocr-v3",
        "version": "3.1.0",
        "moteur_principal":  "paddleocr",
        "moteur_correction": "groq-vision (llama-4-scout)" if _groq else "désactivé (GROQ_API_KEY manquante)",
        "paddle_available":  PADDLE_AVAILABLE,
        "groq_available":    _groq,
        "cv2_available":     CV2_AVAILABLE,
        "docx_disponible":   DOCX_AVAILABLE,
        "xlsx_disponible":   XLSX_AVAILABLE,
        "pptx_disponible":   PPTX_AVAILABLE,
        "formats_supportes": FORMATS_SUPPORTES,
        "capacites": {
            "texte":      True,
            "layout":     CV2_AVAILABLE,
            "figures":    CV2_AVAILABLE,
            "correction": _groq,
        },
        "avertissement": (
            None if PADDLE_AVAILABLE else
            "PaddleOCR non disponible — vérifiez paddlepaddle==3.3.1 paddleocr==3.5.0 "
            "et relancez uvicorn SANS --reload"
        ),
    }
